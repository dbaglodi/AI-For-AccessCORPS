import io
import logging
import numpy as np
from PIL import Image
import torch
from transformers import (
    TableTransformerForObjectDetection,
    DetrImageProcessor,
)
import easyocr
import numpy as np


from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx.util import Pt

logger = logging.getLogger(__name__)

DEVICE = "cuda" if torch.cuda.is_available() else "cpu"
_table_models = None


def get_table_models():
    """Lazy-load Table Transformer models."""
    global _table_models
    if _table_models is not None:
        return _table_models

    logger.info(f"Loading Table Transformer models on {DEVICE}...")
    try:
        # Structure recognition model (identifies rows, columns, cells)
        processor = DetrImageProcessor.from_pretrained(
            "microsoft/table-transformer-structure-recognition-v1.1-all",
            size={"height": 800, "width": 800}  # explicit h/w instead of longest_edge
        )
        model = TableTransformerForObjectDetection.from_pretrained(
            "microsoft/table-transformer-structure-recognition-v1.1-all"
        ).to(DEVICE)
        model.eval()

        _table_models = {"processor": processor, "model": model}
        logger.info("✅ Table Transformer loaded successfully.")
        return _table_models
    except Exception as e:
        logger.error(f"Failed to load Table Transformer: {e}")
        return None


_reader = easyocr.Reader(['en'], gpu=True if DEVICE == "cuda" else False)

def _get_cell_text(img: Image.Image, bbox: list[float]) -> str:
    x0, y0, x1, y1 = [int(v) for v in bbox]
    if x1 <= x0 or y1 <= y0: return ""
    
    cell_img = img.crop((x0, y0, x1, y1))
    
    # EasyOCR expects a numpy array, not a PIL Image
    cell_np = np.array(cell_img)
    
    # Extract text (detail=0 returns just the strings, not the coordinates)
    results = _reader.readtext(cell_np, detail=0)
    
    return " ".join(results).strip()


def _boxes_to_grid(
    row_boxes: list, col_boxes: list, cell_boxes: list, img: Image.Image
) -> list[list[str]]:
    """
    Given sorted row/column bounding boxes, assign each detected cell
    to a (row, col) slot and OCR the text inside.
    """
    # Sort rows top-to-bottom, columns left-to-right
    row_boxes = sorted(row_boxes, key=lambda b: b[1])
    col_boxes = sorted(col_boxes, key=lambda b: b[0])

    n_rows, n_cols = len(row_boxes), len(col_boxes)
    if n_rows == 0 or n_cols == 0:
        return []

    grid = [["" for _ in range(n_cols)] for _ in range(n_rows)]

    def center(box):
        return ((box[0] + box[2]) / 2, (box[1] + box[3]) / 2)

    def nearest_idx(centers, val):
        return int(np.argmin([abs(c - val) for c in centers]))

    row_centers = [(b[1] + b[3]) / 2 for b in row_boxes]
    col_centers = [(b[0] + b[2]) / 2 for b in col_boxes]

    for cell_box in cell_boxes:
        cx, cy = center(cell_box)
        r = nearest_idx(row_centers, cy)
        c = nearest_idx(col_centers, cx)
        if grid[r][c] == "":  # avoid overwriting with duplicates
            grid[r][c] = _get_cell_text(img, cell_box)

    return grid


def extract_table_from_image(image_bytes: bytes, provider: str = "local", api_key: str = None) -> list[list[str]]:
    if provider.lower() == "gemini":
        logger.info("Routing Table extraction directly to Gemini...")
        return _extract_table_with_gemini(image_bytes, api_key)

    # 2. Local Flow
    logger.info("Routing Table extraction to Local (Table Transformer + EasyOCR)...")
    models = get_table_models()
    if not models:
        return []

    try:
        img_original = Image.open(io.BytesIO(image_bytes)).convert("RGB")
        
        # Explicitly resize to match what the processor uses internally
        img_resized = img_original.resize((800, 800), Image.Resampling.LANCZOS)
        
        # Run detection on the resized image
        inputs = models["processor"](images=img_resized, return_tensors="pt").to(DEVICE)
        
        with torch.no_grad():
            outputs = models["model"](**inputs)

        # target_sizes must match what was actually fed to the model
        target_sizes = torch.tensor([[800, 800]], device=DEVICE)
        results = models["processor"].post_process_object_detection(
            outputs, threshold=0.7, target_sizes=target_sizes
        )[0]

        id2label = models["model"].config.id2label
        row_boxes, col_boxes, cell_boxes = [], [], []
        for score, label, box in zip(results["scores"], results["labels"], results["boxes"]):
            label_name = id2label[label.item()].lower()
            coords = box.tolist()
            if "row" in label_name and "header" not in label_name:
                row_boxes.append(coords)
            elif "column" in label_name:
                col_boxes.append(coords)
            elif "cell" in label_name:
                cell_boxes.append(coords)

        logger.info(f"Detected {len(row_boxes)} rows, {len(col_boxes)} cols, {len(cell_boxes)} cells")

        if not cell_boxes and row_boxes and col_boxes:
            for r_box in row_boxes:
                for c_box in col_boxes:
                    x0 = max(r_box[0], c_box[0])
                    y0 = max(r_box[1], c_box[1])
                    x1 = min(r_box[2], c_box[2])
                    y1 = min(r_box[3], c_box[3])
                    if x1 > x0 and y1 > y0:
                        cell_boxes.append([x0, y0, x1, y1])

       # NEW: Calculate scaling factors
        orig_w, orig_h = img_original.size
        scale_x = orig_w / 800.0
        scale_y = orig_h / 800.0

        # NEW: Scale boxes back to original aspect ratio
        row_boxes = [[b[0]*scale_x, b[1]*scale_y, b[2]*scale_x, b[3]*scale_y] for b in row_boxes]
        col_boxes = [[b[0]*scale_x, b[1]*scale_y, b[2]*scale_x, b[3]*scale_y] for b in col_boxes]
        cell_boxes = [[b[0]*scale_x, b[1]*scale_y, b[2]*scale_x, b[3]*scale_y] for b in cell_boxes]

        # Pass img_original instead of img_resized for OCR quality
        grid = _boxes_to_grid(row_boxes, col_boxes, cell_boxes, img_original)
        
        # --- NEW DYNAMIC FALLBACK ROUTING ---
        if not grid or not any(grid):
            logger.warning("Local standard models failed to extract table structure. Returning empty.")
            return []

        return grid

    except Exception as e:
        logger.error(f"Table extraction failed: {e}")
        return []

def _extract_table_with_gemini(image_bytes: bytes, api_key: str) -> list[list[str]]:
    """
    Fallback for borderless/complex tables using Google's Gemini model.
    """
    try:
        import google.generativeai as genai
        import json
        
        if not api_key:
            logger.error("No API key provided for Gemini fallback.")
            return []

        genai.configure(api_key=api_key)
        # Using 1.5-flash as it is extremely fast and great at visual JSON extraction
        model = genai.GenerativeModel('gemini-2.5-flash-lite')

        img = Image.open(io.BytesIO(image_bytes)).convert("RGB")

        prompt = (
            "Extract the table from this image. "
            "Return ONLY a JSON array of arrays where each inner array is a row "
            "and each string is a cell value. "
            "No markdown, no backticks, just raw JSON."
        )

        response = model.generate_content([prompt, img])
        raw = response.text.strip()

        # Clean up common model output artifacts
        raw = raw.replace("```json", "").replace("```", "").strip()

        table_data = json.loads(raw)
        if isinstance(table_data, list) and all(isinstance(r, list) for r in table_data):
            return table_data
        return []

    except Exception as e:
        logger.error(f"Gemini table fallback failed: {e}")
        return []

def insert_table_into_pptx(slide, image_shape, table_data: list[list[str]], alt_text: str = ""):
    if not table_data or not table_data[0]:
        return

    rows = len(table_data)
    cols = max(len(row) for row in table_data)
    x, y = image_shape.left, image_shape.top
    width, height = image_shape.width, image_shape.height
    
    table_shape = slide.shapes.add_table(rows, cols, x, y, width, height)
    table = table_shape.table

    for r, row_data in enumerate(table_data):
        for c, cell_text in enumerate(row_data):
            if c < cols:
                cell = table.cell(r, c)
                cell.text = str(cell_text)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(10)

    # --- NEW: Apply Alt Text to the PowerPoint Table ---
    if alt_text:
        table_shape._element.nvGraphicFramePr.cNvPr.set('descr', alt_text)
    # ---------------------------------------------------

    spTree = slide.shapes._spTree
    spTree.insert(spTree.index(image_shape._element), table_shape._element)

def insert_table_into_docx(doc, shape_element, table_data: list[list[str]], alt_text: str = ""):
    if not table_data or not table_data[0]:
        return

    # Use XPath to directly find the paragraph containing this image
    parent_p = shape_element.xpath('./ancestor::w:p')
    if not parent_p:
        return
        
    target_xml_p = parent_p[0]

    rows = len(table_data)
    cols = max(len(row) for row in table_data)
    table = doc.add_table(rows=rows, cols=cols)
    table.style = 'Table Grid'
    
    for r, row_data in enumerate(table_data):
        for c, cell_text in enumerate(row_data):
            if c < cols:
                table.cell(r, c).text = str(cell_text)
                
    # Apply Alt Text to the Word Table
    if alt_text:
        tblDescr = OxmlElement('w:tblDescription')
        tblDescr.set(qn('w:val'), alt_text)
        table._tbl.tblPr.append(tblDescr)
        
    # Insert the table immediately after the paragraph containing the image
    target_xml_p.addnext(table._tbl)