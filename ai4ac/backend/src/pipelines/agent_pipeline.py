import sys
import os
import shutil
import subprocess
import tempfile
import io
import logging
import time
import json
import re
import threading
import concurrent.futures
from pathlib import Path
from typing import Optional, Dict, Any, List, Tuple

from PIL import Image
from docx import Document
from docx.text.paragraph import Paragraph
from docx.oxml.shape import CT_Inline 
from docx.oxml.ns import nsdecls, qn 
from pptx import Presentation
from pptx.slide import Slide as PptxSlide 
from pptx.shapes.picture import Picture as PptxPicture
from pptx.enum.shapes import MSO_SHAPE_TYPE

logger = logging.getLogger(__name__)
try:
    import torch
    import numpy as np
    from transformers import AutoProcessor, PaliGemmaProcessor, PaliGemmaForConditionalGeneration
except ImportError:
    torch = None
    np = None
    AutoProcessor = None
    PaliGemmaProcessor = None
    PaliGemmaForConditionalGeneration = None
    logger.warning("Heavy ML libraries not found. Local models disabled. Gemini API required.")
import google.generativeai as genai

# Local imports
from src.pipelines.equation_pipeline import extract_equations_from_image, insert_equation_into_docx
from src.pipelines.table_pipeline import extract_table_from_image, insert_table_into_docx, insert_table_into_pptx
# Calculate project root (ai4ac/backend) relative to this file
PROJECT_ROOT = Path(__file__).parent.parent.parent.absolute()
sys.path.append(str(PROJECT_ROOT))

from src.config.gpu_config import get_gpu_settings
from src.models.vision_processor import get_model as get_fallback_model 
from src.config.app_config import CUSTOM_CACHE_DIR

# LangChain availability check
try:
    from langchain.text_splitter import RecursiveCharacterTextSplitter
    from langchain_community.vectorstores import FAISS
    from langchain_huggingface import HuggingFaceEmbeddings
    LANGCHAIN_AVAILABLE = True
except ImportError:
    RecursiveCharacterTextSplitter, FAISS, HuggingFaceEmbeddings = None, None, None
    LANGCHAIN_AVAILABLE = False

# --- Utilities ---

def chunk_text(text: str, chunk_size: int = 512, overlap: int = 64):
    if not text: return []
    if RecursiveCharacterTextSplitter:
        splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=overlap)
        return splitter.split_text(text)
    tokens = text.split(); chunks = []
    for i in range(0, len(tokens), chunk_size - overlap):
        chunks.append(" ".join(tokens[i:i + chunk_size]))
    return chunks

class LangChainRAG:
    def __init__(self, model_name='all-MiniLM-L6-v2'):
        if not LANGCHAIN_AVAILABLE:
            raise ImportError("LangChain components not installed.")
        emb_cache_folder = os.path.join(CUSTOM_CACHE_DIR, 'sentence-transformers')
        self.embeddings = HuggingFaceEmbeddings(model_name=model_name, cache_folder=emb_cache_folder)
        self.vector_store = None
    def add_documents(self, texts):
        if not texts: return
        self.vector_store = FAISS.from_texts(texts, self.embeddings)
    def search(self, query, top_k=3):
        if not self.vector_store: return []
        results = self.vector_store.similarity_search_with_score(query, k=top_k)
        return [{'text': doc.page_content, 'score': score} for doc, score in results]

class SimpleTextRAG:
    def __init__(self): self.chunks = []
    def add_documents(self, texts, metadatas=None): self.chunks.extend(texts)
    def search(self, query, top_k=4):
        if not self.chunks: return []
        query_words = set(query.lower().split())
        results = [{'text': chunk, 'score': len(query_words.intersection(set(chunk.lower().split())))} for chunk in self.chunks]
        results.sort(key=lambda x: x['score'], reverse=True)
        return [r for r in results if r['score'] > 1][:top_k]

# Global settings and state
MAX_CHUNKS = int(os.environ.get("AGENT_MAX_CHUNKS", 1000))
DEFAULT_MAX_WORKERS = min(4, (os.cpu_count() or 1))
_primary_model_cache = None
_primary_model_lock = threading.Lock()
_results_lock = threading.Lock()

UNSUPPORTED_IMAGE_PLACEHOLDER = "https://placehold.co/400x300/EFEFEF/AAAAAA?text=Unsupported+Format%5Cn(WMF/EMF)"

# --- Image Processing ---

def try_convert_metafile(image_bytes: bytes) -> Optional[bytes]:
    """Attempts to convert WMF/EMF to PNG using system tools (magick, convert, or libreoffice)."""
    input_path = None
    output_path = None
    
    local_magick = str(PROJECT_ROOT / "magick")
    
    try:
        with tempfile.NamedTemporaryFile(suffix=".wmf", delete=False) as f_in:
            f_in.write(image_bytes)
            input_path = f_in.name
        
        output_path = input_path.replace(".wmf", ".png")
        
        # 1. Try ImageMagick commands
        for cmd in [local_magick, 'magick', 'convert']:
            try:
                # If path contains slashes, check existence first
                if "/" in cmd and not os.path.exists(cmd):
                    continue
                    
                actual_args = [cmd, 'convert', input_path, output_path] if "magick" in cmd else [cmd, input_path, output_path]
                
                result = subprocess.run(actual_args, capture_output=True, timeout=10)
                if result.returncode == 0 and os.path.exists(output_path):
                    with open(output_path, "rb") as f_out:
                        return f_out.read()
            except (FileNotFoundError, PermissionError):
                continue

        # 2. Fallback: Try LibreOffice
        try:
            subprocess.run([
                'libreoffice', '--headless', '--convert-to', 'png', 
                input_path, '--outdir', os.path.dirname(input_path)
            ], capture_output=True, timeout=15)
            
            base = os.path.splitext(os.path.basename(input_path))[0]
            lo_output = os.path.join(os.path.dirname(input_path), f"{base}.png")
            if os.path.exists(lo_output):
                output_path = lo_output 
                with open(lo_output, "rb") as f_out:
                    return f_out.read()
        except FileNotFoundError:
             pass 

    except Exception as e:
        logger.warning(f"Metafile conversion failed: {e}")
    finally:
        try:
            if input_path and os.path.exists(input_path): os.remove(input_path)
            if output_path and os.path.exists(output_path): os.remove(output_path)
        except: pass
    return None

def preprocess_image(image_bytes: bytes, max_size: int = 512) -> bytes:
    try:
        try:
            image = Image.open(io.BytesIO(image_bytes)).convert('RGB')
        except Exception:
            converted = try_convert_metafile(image_bytes)
            if converted:
                image = Image.open(io.BytesIO(converted)).convert('RGB')
            else:
                raise ValueError("Unsupported image format")

        image.thumbnail((max_size, max_size), Image.Resampling.LANCZOS)
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='JPEG', quality=90)
        return img_byte_arr.getvalue()
    except Exception as e:
        if str(e) != "Unsupported image format":
            logger.error(f"Error preparing image for prompt: {e}")
        raise e

# --- Model Loading ---
def get_primary_model(provider="local", logger_instance=None):
    global _primary_model_cache
    log = logger_instance or logger

    # 1. LAZY LOADING INTERCEPT: If Gemini, return None immediately without locking or loading!
    if provider == "gemini":
        log.info("Gemini provider selected. Bypassing local model load.")
        return None 

    # 2. Local Model Loading (with existing thread lock)
    with _primary_model_lock:
        # If it was already loaded by a previous request, return it instantly
        if _primary_model_cache: 
            return _primary_model_cache
            
        # Otherwise, this is the Cold Start
        log.info("First request for local model detected. Downloading/Loading weights into VRAM now...")
        gpu_settings = get_gpu_settings()
        primary_model_id = "google/paligemma-3b-mix-448"
        try:
            log.info(f"Loading primary vision model: {primary_model_id}")
            processor = PaliGemmaProcessor.from_pretrained(primary_model_id, cache_dir=CUSTOM_CACHE_DIR)
            model = PaliGemmaForConditionalGeneration.from_pretrained(
                primary_model_id,
                torch_dtype=gpu_settings["dtype"],
                attn_implementation="eager",
                cache_dir=CUSTOM_CACHE_DIR
            ).to(gpu_settings["device"])
            
            _primary_model_cache = {
                "model": model, 
                "processor": processor,
                "text_rag": LangChainRAG() if LANGCHAIN_AVAILABLE else SimpleTextRAG(),
                "type": "vision_model_paligemma"
            }
            return _primary_model_cache
            
        except Exception as e:
            log.error(f"Primary model load failed: {e}")
            _primary_model_cache = {"model": None, "processor": None, "text_rag": SimpleTextRAG(), "type": "failed"}
            return _primary_model_cache

# --- Prompts ---

def _get_combined_context(structured_context: Dict[str, Optional[str]]) -> str:
    parts = []
    if structured_context.get("doc_title"): parts.append(f"Document: {structured_context['doc_title']}")
    if structured_context.get("slide_title"): parts.append(f"Title: {structured_context['slide_title']}")
    if structured_context.get("surrounding_text"): parts.append(f"Immediate Text: {structured_context['surrounding_text'][:300]}")
    # --- ADDED RAG CONTEXT ---
    if structured_context.get("broader_context"): parts.append(f"Broader Document Context: {structured_context['broader_context'][:500]}")
    
    return ". ".join(filter(None, parts)) or "No context available."

def create_tagging_prompt(structured_context, model_format="paligemma"):
    ctx = _get_combined_context(structured_context)
    prompt_base = "Categories: Graph, Map, Diagram, Table, Photograph, Text, Screenshot, Equation, Other. Respond ONLY with comma-separated names of applicable categories, RANKED by relevance."
    if model_format == "paligemma":
        return f"<image>\nContext: {ctx}\n{prompt_base}\nSelected Categories:"
    return f"<image>\nContext: {ctx}\nQuestion: {prompt_base}\nAnswer:"

def create_complex_data_alt_text_prompt(structured_context, categories, existing_alt, model_format="paligemma"):
    ctx = _get_combined_context(structured_context)
    cat_str = ", ".join(categories)
    instr = f"Image Type: {cat_str}. 1. Identify main trend. 2. Describe axes. 3. DO NOT list individual points. 4. CRITICAL: Do not read labels literally; explain context. 5. Keep the description under 300 characters"
    p = f"Context: {ctx}\nExisting Alt: {existing_alt}\n{instr}\nAlt Text:"
    return f"<image>\n{p}" if model_format == "paligemma" else f"<image>\n{p}\nAnswer:"

def create_alt_text_prompt(structured_context, categories, existing_alt, model_format="paligemma"):
    ctx = _get_combined_context(structured_context)
    cat_str = ", ".join(categories)
    instr = f"Describe this {cat_str} focusing on visual elements and context relationship."
    if model_format == "gemini":
        instr += """
        Important guidelines:
    1. Focus on describing key visual elements and their relationship to the context
    2. Use clear, academic language appropriate for the content
    3. Keep the description under 125 characters
    4. If the image shows a diagram or figure, describe its key components and purpose
    5. Include relevant technical terms from the context when appropriate
    6. Consider how this image fits into the overall presentation theme
    """
    p = f"Context: {ctx}\nExisting Alt: {existing_alt}\n{instr}\nAlt Text:"
    return f"<image>\n{p}" if model_format == "paligemma" else f"<image>\n{p}\nAnswer:"

def create_mathml_prompt(model_format="paligemma"):
    prompt = "Identify the mathematical equation in the image and convert it to valid MathML 3.0. Return only the MathML code."
    return f"<image>\n{prompt}\nMathML:" if model_format == "paligemma" else f"<image>\nQuestion: {prompt}\nAnswer:"

# --- Alt Text Pipeline ---

def classify_and_generate_alt_text(
    image_bytes, structured_context, primary_model_system=None, 
    ext=".pptx", existing_alt="", slide_num=None, task_idx=0,
    provider="local", api_key=None, forced_pipeline=None # <--- ADD THIS HERE
):
    try:
        prompt_image_bytes = preprocess_image(image_bytes)
    except Exception as e:
        return ["Needs Review"], f"Preprocessing error: {e}", ""

    alt_text = ""
    mathml = ""
    categories = ["Other"]
    complex_types = {"Graph", "Map", "Diagram"}
    valid_map = {c.lower(): c for c in ["Graph", "Map", "Diagram", "Figure", "Table", "Photograph", "Text", "Screenshot", "Equation", "Other"]}

    gemini_success = False

    # --- NEW GEMINI BRANCH WITH 3 RETRIES ---
    if provider == "gemini" and api_key:
        max_retries = 3
        for attempt in range(max_retries):
            try:
                genai.configure(api_key=api_key)
                model = genai.GenerativeModel('gemini-2.5-flash-lite')
                image_pil = Image.open(io.BytesIO(prompt_image_bytes))

                # Pass 1: Tagging or Override
                if forced_pipeline and forced_pipeline.lower() != "general":
                    categories = [forced_pipeline]
                    logger.info(f"Bypassing Gemini classification. User forced pipeline: {forced_pipeline}")
                else:
                    p1 = create_tagging_prompt(structured_context, "gemini").replace("<image>\n", "")
                    response1 = model.generate_content([p1, image_pil])
                    tag_text = re.sub(r"^(Category:|Answer:|Selected Categories:)\s*", "", response1.text, flags=re.IGNORECASE).strip()
                    categories = [valid_map[t.strip().lower()] for t in tag_text.split(',') if t.strip().lower() in valid_map] or ["Other"]

                # Pass 2: Alt Text
                is_complex = any(c in complex_types for c in categories)
                p2 = create_complex_data_alt_text_prompt(structured_context, categories, existing_alt, "gemini") if is_complex else create_alt_text_prompt(structured_context, categories, existing_alt, "gemini")
                p2 = p2.replace("<image>\n", "")
                response2 = model.generate_content([p2, image_pil])
                alt_text = response2.text

                # Pass 3: MathML
                if "Equation" in categories:
                    p3 = create_mathml_prompt("gemini").replace("<image>\n", "")
                    response3 = model.generate_content([p3, image_pil])
                    mathml = response3.text

                gemini_success = True
                break # Success! Exit the retry loop early

            except Exception as e:
                logger.warning(f"Gemini API attempt {attempt + 1} failed for image {task_idx}: {e}")
                if attempt < max_retries - 1:
                    logger.info(f"Waiting 30 seconds before retrying Gemini for image {task_idx}...")
                    time.sleep(30)
                else:
                    logger.error(f"All {max_retries} Gemini attempts failed. Falling back to local model.")
                    # Keep alt_text empty so it drops into the local fallback

    # --- EXISTING LOCAL MODEL BRANCH (USED IF LOCAL SELECTED OR IF GEMINI FAILED) ---
    if not gemini_success:
        # PREVENT MEMORY CRASH: Do NOT load the massive local model if Gemini fails due to rate limits.
        if provider == "gemini":
            logger.error("Gemini failed (likely rate limit). Returning error text to prevent server crash.")
            return ["Needs Review"], "API Error: Gemini rate limit reached (15 RPM). Please wait 1 minute and click Regenerate.", ""

        if primary_model_system and primary_model_system.get("model"):
            try:
                model, processor = primary_model_system["model"], primary_model_system["processor"]
                device = get_gpu_settings()["device"]
                image = Image.open(io.BytesIO(prompt_image_bytes))

                # Pass 1: Tagging
                p1 = create_tagging_prompt(structured_context, "paligemma")
                inputs = processor(text=p1, images=image, return_tensors="pt").to(device)
                with torch.inference_mode():
                    out = model.generate(**inputs, max_new_tokens=50)
                    gen_tags = processor.decode(out[0][inputs["input_ids"].shape[1]:], skip_special_tokens=True).strip()
                
                categories = [valid_map[t.strip().lower()] for t in gen_tags.split(',') if t.strip().lower() in valid_map] or ["Other"]
                
                # Pass 2: Branching for Alt Text
                is_complex = any(c in complex_types for c in categories)
                p2 = create_complex_data_alt_text_prompt(structured_context, categories, existing_alt, "paligemma") if is_complex else create_alt_text_prompt(structured_context, categories, existing_alt, "paligemma")
                
                inputs2 = processor(text=p2, images=image, return_tensors="pt").to(device)
                with torch.inference_mode():
                    out2 = model.generate(**inputs2, max_new_tokens=250)
                    alt_text = processor.decode(out2[0][inputs2["input_ids"].shape[1]:], skip_special_tokens=True).strip()

                # Pass 3: MathML conversion if it's an Equation
                if "Equation" in categories:
                    p3 = create_mathml_prompt("paligemma")
                    inputs3 = processor(text=p3, images=image, return_tensors="pt").to(device)
                    with torch.inference_mode():
                        out3 = model.generate(**inputs3, max_new_tokens=500)
                        mathml = processor.decode(out3[0][inputs3["input_ids"].shape[1]:], skip_special_tokens=True).strip()

            except Exception as e:
                logger.error(f"Primary model prompt failed for image {task_idx}: {e}")

        # Fallback (SmolVLM)
        if not alt_text:
            try:
                fb_info = get_fallback_model()
                if fb_info and fb_info.get("model"):
                    from src.models.vision_processor import process_image as fb_process
                    
                    # Tagging Fallback
                    p1_fb = create_tagging_prompt(structured_context, "smolvlm")
                    tag_fb = fb_process(prompt_image_bytes, p1_fb)
                    tag_fb = re.sub(r"^(Category:|Answer:|Selected Categories:)\s*", "", tag_fb, flags=re.IGNORECASE).strip()
                    categories = [valid_map[t.strip().lower()] for t in tag_fb.split(',') if t.strip().lower() in valid_map] or ["Other"]
                    
                    # Alt Text Fallback
                    is_complex = any(c in complex_types for c in categories)
                    p2_fb = create_complex_data_alt_text_prompt(structured_context, categories, existing_alt, "smolvlm") if is_complex else create_alt_text_prompt(structured_context, categories, existing_alt, "smolvlm")
                    alt_text = fb_process(prompt_image_bytes, p2_fb)

                    # MathML Fallback
                    if "Equation" in categories and not mathml:
                        p3_fb = create_mathml_prompt("smolvlm")
                        mathml = fb_process(prompt_image_bytes, p3_fb)
            except Exception as e:
                logger.error(f"Fallback prompt failed for image {task_idx}: {e}")

    # Final Post-processing
    alt_text = re.sub(r"^(Answer:|Alt Text:)\s*", "", alt_text, flags=re.IGNORECASE).strip()
    alt_text = alt_text.replace("<|end|>", "").replace("<image>", "").strip()
    
    mathml = re.sub(r"^(Answer:|MathML:)\s*", "", mathml, flags=re.IGNORECASE).strip()
    mathml = mathml.replace("<|end|>", "").replace("<image>", "").strip()

    if any(p in alt_text for p in ["/ba/", "/da/", "/ga/"]):
        logger.warning(f"Problematic text reading detected in image {task_idx}. Clearing.")
        alt_text = ""; categories.append("Needs Review")

    if not alt_text:
        alt_text = existing_alt or f"Image {task_idx+1}"
        if "Needs Review" not in categories: categories.append("Needs Review")

    return categories, alt_text, mathml

# --- Context Extraction ---

def get_context_for_image_pptx(slide: PptxSlide, shape: PptxPicture) -> Dict[str, Optional[str]]:
    ctx = {"slide_title": None, "surrounding_text": ""}
    try:
        if slide.shapes.title and slide.shapes.title.text.strip():
            ctx["slide_title"] = slide.shapes.title.text.strip().split('\n')[0]
        txt = [s.text.strip() for s in slide.shapes if hasattr(s, "text") and s.text.strip()]
        ctx["surrounding_text"] = " ".join(txt[:2])
    except: pass
    return ctx

def get_context_for_image_docx(doc: Document, inline_shape) -> Dict[str, Optional[str]]:
    ctx = {"doc_title": getattr(doc.core_properties, 'title', None), "surrounding_text": ""}
    try:
        for p in doc.paragraphs:
            # Look for the specific inline shape in the paragraph XML
            if inline_shape in p._element.xpath('.//wp:inline'):
                ctx["surrounding_text"] = p.text[:300]
                break
    except: pass
    return ctx

# --- Main Entry Point ---

def run_agent_pipeline(file_path, ext, progress_callback=None, provider="local", api_key=None, **kwargs):
    gpu_settings = get_gpu_settings()
    primary_model = get_primary_model(provider=provider)
    results = []
    
    # 1. --- BACKGROUND RAG INGESTION ---
    text_rag = primary_model.get("text_rag") if primary_model else None
    if text_rag:
        doc_texts = []
        try:
            if ext == ".docx":
                doc_temp = Document(file_path)
                for p in doc_temp.paragraphs:
                    if p.text.strip(): doc_texts.append(p.text.strip())
            elif ext == ".pptx":
                pres_temp = Presentation(file_path)
                for slide in pres_temp.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            doc_texts.append(shape.text.strip())
            
            # Chunk and ingest
            chunks = []
            for text in doc_texts:
                chunks.extend(chunk_text(text))
            if chunks:
                text_rag.add_documents(chunks)
                logger.info(f"Ingested {len(chunks)} text chunks into RAG.")
        except Exception as e:
            logger.warning(f"Background RAG ingestion failed: {e}")
    
    try:
        if ext == ".pptx":
            pres = Presentation(file_path)
            images = []
            pres_modified = False 
            
            # 1. Use the expanded list to find all visual data
            valid_types = [
                MSO_SHAPE_TYPE.PICTURE, 
                MSO_SHAPE_TYPE.LINKED_PICTURE,
                MSO_SHAPE_TYPE.CHART, 
                MSO_SHAPE_TYPE.IGX_GRAPHIC, # SmartArt
                MSO_SHAPE_TYPE.GROUP
            ]

            for slide_num, slide in enumerate(pres.slides, 1):
                for shape in slide.shapes:
                    if getattr(shape, "shape_type", None) in valid_types:
                        images.append((slide_num, shape))
            
            total_count = len(images)
            for i, (slide_num, shape) in enumerate(images, 1):
                try:
                    # 2. Safely get existing alt text for any shape type using XPath
                    alt = ""
                    cNvPr_elements = shape._element.xpath('.//*[local-name()="cNvPr"]')
                    if cNvPr_elements:
                        alt = cNvPr_elements[0].get('descr', '')

                    ctx = get_context_for_image_pptx(pres.slides[slide_num-1], shape)
                    
                    # 3. Safely extract image bytes (Handle Charts/Groups)
                    image_bytes = None
                    if getattr(shape, "shape_type", None) in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE]:
                        try:
                            if hasattr(shape, "image"): image_bytes = shape.image.blob
                        except AttributeError: pass

                    if image_bytes is None:
                        # FALLBACK: Port your Colab's screenshot-and-crop logic here!
                        # Because this is the backend, if you haven't brought over the 
                        # pdf2image / libreoffice slide screenshot function yet, 
                        # we must log a warning and skip to avoid crashing the server.
                        logger.warning(f"Shape {i} (Type {shape.shape_type}) requires slide cropping. Ensure screenshot function is implemented.")
                        continue # Remove this continue once your slide screenshot function is ported

                    # --- BACKGROUND RAG QUERY (PPTX) ---
                    if text_rag and ctx.get("surrounding_text"):
                        rag_hits = text_rag.search(ctx["surrounding_text"], top_k=2)
                        if rag_hits:
                            ctx["broader_context"] = " ".join([hit['text'] for hit in rag_hits])
                    
                    # Call Gemini / Local Model
                    cats, gen_alt, _ = classify_and_generate_alt_text(
                        image_bytes, ctx, primary_model, ext, alt, slide_num, i-1,
                        provider=provider, api_key=api_key
                    )
                    
                    # --- TABLE PIPELINE INTEGRATION ---
                    if "Table" in cats: # <-- Relaxed condition
                        logger.info(f"Image {i} classification includes Table. Extracting data...")
                        # Pass provider and api_key here
                        table_data = extract_table_from_image(shape.image.blob, provider=provider, api_key=api_key)
                        if table_data:
                            insert_table_into_pptx(pres.slides[slide_num-1], shape, table_data, alt_text=gen_alt)
                            pres_modified = True

                    import base64
                    try:
                        disp_bytes = preprocess_image(shape.image.blob, 256)
                    except:
                        disp_bytes = shape.image.blob # Best effort
                    
                    img_data = base64.b64encode(disp_bytes).decode()
                    results.append({"classification": cats, "alt_text": alt, "generated_alt_text": gen_alt, "image_idx": i, "slide_num": slide_num, "image_data": f"data:image/jpeg;base64,{img_data}"})
                except Exception as inner_e:
                    logger.error(f"Error on pptx image {i}: {inner_e}")
                    results.append({"classification": ["Needs Review"], "alt_text": "", "generated_alt_text": f"Error: {inner_e}", "image_idx": i, "slide_num": slide_num})
                
                if progress_callback: progress_callback(f"Image {i}/{total_count}", i, total_count)
                
            if pres_modified:
                pres.save(file_path)

        elif ext == ".docx":
            doc = Document(file_path)
            
            # Use XPath to find ALL images (inline and anchored)
            drawing_elements = doc._element.xpath('.//w:drawing')
            total_count = len(drawing_elements)
            doc_modified = False 
            
            for i, drawing in enumerate(drawing_elements, 1):
                try:
                    # Safely get Relationship ID
                    # REMOVED: namespaces=namespaces (python-docx handles this internally)
                    rId = None
                    blips = drawing.xpath('.//a:blip/@r:embed')
                    if blips:
                        rId = blips[0]
                    else:
                        charts = drawing.xpath('.//c:chart/@r:id')
                        if charts: rId = charts[0]

                    if not rId: continue
                    
                    rel = doc.part.rels[rId]
                    image_bytes = rel.target_part.blob
                    
                    # Safely get alt text
                    alt = ""
                    docPr_elements = drawing.xpath('.//wp:docPr')
                    if docPr_elements:
                        alt = docPr_elements[0].get('descr', '')

                    # Get surrounding context
                    ctx = {"surrounding_text": ""}
                    parent_p = drawing.xpath('./ancestor::w:p')
                    if parent_p:
                        ctx["surrounding_text"] = "".join(parent_p[0].itertext())[:300]
                    
                    # --- BACKGROUND RAG QUERY (DOCX) ---
                    if text_rag and ctx.get("surrounding_text"):
                        rag_hits = text_rag.search(ctx["surrounding_text"], top_k=2)
                        if rag_hits:
                            ctx["broader_context"] = " ".join([hit['text'] for hit in rag_hits])
                            
                    cats, gen_alt, _ = classify_and_generate_alt_text(
                        image_bytes, ctx, primary_model, ext, alt, 1, i-1,
                        provider=provider, api_key=api_key
                    )
                    
                    # --- TABLE PIPELINE INTEGRATION ---
                    if "Table" in cats:
                        logger.info(f"Image {i} classification includes Table. Extracting data...")
                        # Pass provider and api_key here
                        table_data = extract_table_from_image(image_bytes, provider=provider, api_key=api_key)
                        if table_data:
                            insert_table_into_docx(doc, drawing, table_data, alt_text=gen_alt)
                            doc_modified = True
                    # ----------------------------------
                    
                    # --- EQUATION PIPELINE INTEGRATION ---
                    if "Equation" in cats:
                        logger.info(f"Image {i} classification includes Equation. Extracting LaTeX...")
                        # Pass provider and api_key here
                        equations = extract_equations_from_image(image_bytes, provider=provider, api_key=api_key)
                        if equations:
                            insert_equation_into_docx(doc, drawing, equations)
                            doc_modified = True
                    # ----------------------------------

                    # Generate preview
                    import base64
                    try: disp_bytes = preprocess_image(image_bytes, 256)
                    except: disp_bytes = image_bytes
                        
                    img_data = base64.b64encode(disp_bytes).decode()
                    results.append({
                        "classification": cats, "alt_text": alt, 
                        "generated_alt_text": gen_alt, "image_idx": i, 
                        "image_data": f"data:image/jpeg;base64,{img_data}", "rId": rId
                    })
                    
                except Exception as inner_e:
                    logger.error(f"Error on docx image {i}: {inner_e}")
                    results.append({"classification": ["Needs Review"], "alt_text": "", "generated_alt_text": f"Error: {inner_e}", "image_idx": i})
                
                if progress_callback: progress_callback(f"Image {i}/{total_count}", i, total_count)
            
            if doc_modified:
                doc.save(file_path)
        
        return results
    except Exception as e:
        logger.error(f"Pipeline crashed: {e}")
        return []