from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Body, BackgroundTasks
from fastapi.responses import FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
import shutil
import glob
import os
from uuid import uuid4
import json
from typing import Dict, List, Optional, Any, Callable
import importlib
import asyncio
from datetime import datetime
import logging
from dotenv import load_dotenv
import sys
from pathlib import Path
from pydantic import BaseModel
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Add project root to Python path
project_root = Path(__file__).parent.parent
sys.path.insert(0, str(project_root))

# --- START MODIFICATION: Import config paths ---
from src.config.app_config import UPLOAD_DIR, PROCESSED_DIR, REMEDIATED_DIR, CUSTOM_CACHE_DIR
# --- END MODIFICATION ---

# --- START MODIFICATION: Import pptx/docx classes ---
from docx import Document
from pptx import Presentation
from docx.oxml.ns import qn
from pptx.shapes.picture import Picture as PptxPicture
# --- END MODIFICATION ---

app = FastAPI()

class RegenerateRequest(BaseModel):
    image_idx: int
    forced_pipeline: str
    slide_num: Optional[int] = None
    rId: Optional[str] = None
    provider: Optional[str] = "local"
    api_key: Optional[str] = None

@app.on_event("startup")
def startup_event():
    """Check disk space on startup."""
    try:
        threshold_gb = 15.0
        cache_dir = CUSTOM_CACHE_DIR
        logging.warning(f"Using custom cache directory for disk space check: {cache_dir}")
        if os.path.exists(cache_dir):
            total, used, free = shutil.disk_usage(str(cache_dir))
            free_gb = free / (1024**3)
            logging.warning(f"Checking Hugging Face cache disk space: {free_gb:.2f} GB free.")
            if free_gb < threshold_gb:
                logging.warning(f"Free space below {threshold_gb} GB. Clearing cache...")
                shutil.rmtree(cache_dir)
                logging.warning("Cache cleared.")
                os.makedirs(cache_dir, exist_ok=True)
        else:
            logging.warning(f"Cache directory not found: {cache_dir}. Skipping check.")
    except Exception as e:
        logging.error(f"Error during startup cache check: {e}")

load_dotenv()

# --- START MODIFICATION: Expose Content-Disposition Header ---
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], # In production, restrict this to your frontend's origin
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["Content-Disposition"] # Allow frontend to read this header
)
# --- END MODIFICATION ---

processing_status = {}

def _resolve_agent_pipeline():
    # (function remains the same)
    try:
        mod = importlib.import_module("src.pipelines.agent_pipeline")
        logging.info("Using agent_pipeline implementation")
        return mod.run_agent_pipeline
    except Exception as exc:
        logging.error(f"Failed to import pipelines.agent_pipeline as agent_pipeline module: {exc}")
        raise

run_agent_pipeline = _resolve_agent_pipeline()

_simple_rag = None
def get_simple_rag():
    # (function remains the same)
    global _simple_rag
    if _simple_rag is not None: return _simple_rag
    try:
        import src.pipelines.agent_pipeline as ap
        if hasattr(ap, 'SimpleRAG'):
            _simple_rag = ap.SimpleRAG(); logging.info('Initialized SimpleRAG from agent_pipeline')
            return _simple_rag
        logging.warning('agent_pipeline does not expose SimpleRAG'); return None
    except Exception as exc: logging.warning(f'Could not initialize SimpleRAG via agent_pipeline: {exc}'); return None

@app.post("/upload")
async def upload_file( 
    background_tasks: BackgroundTasks, 
    file: UploadFile = File(description="Document file to process"),
    provider: str = Form("local"),
    api_key: Optional[str] = Form(None),
    start_slide: Optional[int] = Form(None), # Added
    end_slide: Optional[int] = Form(None)    # Added
):
    # (function remains the same)
    if not file: raise HTTPException(status_code=400, detail="No file uploaded")
    file_id = None
    try:
        ext = os.path.splitext(file.filename)[1].lower()
        if ext not in [".docx", ".pptx"]: raise HTTPException(status_code=400, detail="Only .docx and .pptx files are supported.")
        file_id = str(uuid4()); save_path = os.path.join(UPLOAD_DIR, f"{file_id}{ext}")
        processing_status[file_id] = {"status": "uploading", "progress": 0, "current_step": "Uploading file",
                                      "total_images": 0, "processed_images": 0, "error": None,
                                      "started_at": datetime.now().isoformat(), "updated_at": datetime.now().isoformat(),
                                      "filename": file.filename # Store original filename
                                     }
        try:
            with open(save_path, "wb") as buffer: content = await file.read(); buffer.write(content)
            processing_status[file_id].update({"status": "processing", "progress": 20, "current_step": "Starting analysis",
                                               "updated_at": datetime.now().isoformat()})
        except Exception as e:
            processing_status[file_id].update({"status": "error", "error": str(e), "updated_at": datetime.now().isoformat()})
            raise HTTPException(status_code=500, detail=f"Failed to save file: {str(e)}")
        finally: await file.close()
        background_tasks.add_task(process_document, file_id, save_path, ext, provider, api_key, start_slide, end_slide) # Updated
        return {"file_id": file_id, "filename": file.filename}
    except HTTPException: raise
    except Exception as e:
        if file_id and file_id in processing_status:
            processing_status[file_id].update({"status": "error", "error": str(e), "updated_at": datetime.now().isoformat()})
        raise HTTPException(status_code=500, detail=f"Unexpected error: {str(e)}")

def process_document(file_id: str, save_path: str, ext: str, provider: str = "local", api_key: Optional[str] = None, start_slide: Optional[int] = None, end_slide: Optional[int] = None):
    # (function remains the same)
    try:
        processing_status[file_id].update({"current_step": "Extracting images/context", "progress": 30, "updated_at": datetime.now().isoformat()})
        partial_dir = os.path.join(PROCESSED_DIR, file_id); images_partial_dir = os.path.join(partial_dir, "images")
        os.makedirs(images_partial_dir, exist_ok=True)
        rag_strategy = os.environ.get('AGENT_RAG_STRATEGY', 'rag')
        results = run_agent_pipeline( save_path, ext,
            progress_callback=lambda step, prog, total: update_processing_status(file_id, step, prog, total),
            partial_save_dir=partial_dir, rag_strategy=rag_strategy,
            provider=provider, api_key=api_key, 
            start_slide=start_slide, end_slide=end_slide ) # Passed to pipeline
        with open(os.path.join(PROCESSED_DIR, f"{file_id}.json"), "w", encoding="utf-8") as f: json.dump(results, f)
        processing_status[file_id].update({"status": "completed", "progress": 100, "current_step": "Processing complete",
                                           "updated_at": datetime.now().isoformat()})
    except Exception as e:
        logging.exception(f"Error in process_document for file_id {file_id}")
        processing_status[file_id].update({"status": "error", "error": str(e), "updated_at": datetime.now().isoformat()})
        if os.path.exists(save_path): os.remove(save_path)

def update_processing_status(file_id: str, step: str, progress: int, total: int):
    # (function remains the same)
    if file_id in processing_status:
        current_total = processing_status[file_id].get("total_images", 0) or 0; new_total = total
        if current_total > 1 and (not new_total or new_total <= 1): new_total = current_total
        processing_status[file_id].update({"current_step": step, "total_images": new_total, "processed_images": progress,
                                           "progress": min(100, 30 + (60 * (progress / max(new_total, 1)))),
                                           "updated_at": datetime.now().isoformat()})
        # Use logging instead of print for status updates if desired
        # logging.info(f"Status update: {step} - {progress}/{new_total} images processed")


@app.get("/status/{file_id}")
def get_status(file_id: str):
    # (function remains the same)
    if file_id not in processing_status: raise HTTPException(status_code=404, detail="File not found")
    return JSONResponse(content=processing_status[file_id])

@app.get("/images/{file_id}")
def get_images(file_id: str):
    # (function remains the same)
    if file_id in processing_status:
        status = processing_status[file_id]
        if status["status"] == "error": raise HTTPException(status_code=500, detail=status["error"])
        elif status["status"] != "completed":
            partial_dir = os.path.join(PROCESSED_DIR, file_id); images_dir = os.path.join(partial_dir, "images")
            partial_images = []
            if os.path.exists(images_dir):
                # Ensure correct sorting if filenames are like 0.json, 1.json, ... 10.json
                files = sorted(glob.glob(os.path.join(images_dir, "*.json")), key=lambda x: int(os.path.basename(x).split('.')[0]))
                for fpath in files:
                    try:
                        with open(fpath, "r", encoding="utf-8") as fh: partial_images.append(json.load(fh))
                    except Exception as e:
                         logging.warning(f"Failed to load partial image data from {fpath}: {e}")
                         continue
            return JSONResponse(content={"status": status["status"], "progress": status["progress"],
                                         "current_step": status["current_step"], "images": partial_images})
    json_path = os.path.join(PROCESSED_DIR, f"{file_id}.json")
    if not os.path.exists(json_path): raise HTTPException(status_code=404, detail="File processing data not found.")
    try:
        with open(json_path, "r", encoding="utf-8") as f: results = json.load(f)
    except Exception as e:
        logging.error(f"Failed to read processed data for {file_id}: {e}")
        raise HTTPException(status_code=500, detail="Failed to read processing results.")
    return JSONResponse(content={"status": "completed", "progress": 100, "current_step": "Processing complete",
                                 "images": results})

@app.post('/rag/ingest')
def rag_ingest(docs: List[Dict[str, Any]]):
    # (function remains the same)
    rag = get_simple_rag();
    if rag is None: raise HTTPException(status_code=503, detail='RAG pipeline not available')
    rag.ingest_documents(docs); return {'status': 'ingested', 'count': len(docs)}

@app.post('/rag/query')
def rag_query(query: Dict[str, Any]):
    # (function remains the same)
    rag = get_simple_rag();
    if rag is None: raise HTTPException(status_code=503, detail='RAG pipeline not available')
    q = query.get('query') or ''; top_k = int(query.get('top_k', 4)); hits = rag.retrieve(q, top_k=top_k)
    return {'query': q, 'hits': hits}

@app.post("/alt-text/{file_id}")
async def update_alt_text( file_id: str, data: Dict[str, Any] = Body(description="Updates to image alt texts") ):
    # (function remains the same)
    json_path = os.path.join(PROCESSED_DIR, f"{file_id}.json")
    if not os.path.exists(json_path): raise HTTPException(status_code=404, detail="File processing data not found")
    try:
        with open(json_path, "r", encoding="utf-8") as f: results = json.load(f)
    except Exception as e:
        logging.error(f"Failed to read processed data for update {file_id}: {e}")
        raise HTTPException(status_code=500, detail="Failed to read processing results.")

    updates = data.get("updates", []); updated_count = 0
    if len(updates) != len(results):
         logging.warning(f"Update/results count mismatch for {file_id} ({len(updates)} vs {len(results)}). Proceeding cautiously.")

    for i, update_data in enumerate(updates):
        if i < len(results):
            # Use the user-provided text if available, otherwise keep the generated one
            final_text = update_data.get("alt_text") # Assuming frontend sends edited text in 'alt_text'
            if final_text is not None: # Check if the key exists, even if the value is empty string
                results[i]["final_alt_text"] = final_text
                updated_count += 1
            elif "final_alt_text" not in results[i]: # If user didn't edit AND it wasn't set before
                results[i]["final_alt_text"] = results[i].get("generated_alt_text", "") # Default to generated

    try:
        with open(json_path, "w", encoding="utf-8") as f: json.dump(results, f)
    except Exception as e:
        logging.error(f"Failed to write updated data for {file_id}: {e}")
        raise HTTPException(status_code=500, detail="Failed to save alt text updates.")

    return {"status": "success", "updated_count": updated_count}

@app.post("/api/regenerate-image/{file_id}")
async def regenerate_image(file_id: str, req: RegenerateRequest):
    """Forces the pipeline to re-run for a specific image using a manually selected category."""
    try:
        # 1. Find the original file
        orig_filename = processing_status.get(file_id, {}).get("filename", "")
        ext = os.path.splitext(orig_filename)[1].lower() if orig_filename else None
        
        orig_path = UPLOAD_DIR / f"{file_id}{ext}"
        if not orig_path.exists():
            found_files = list(UPLOAD_DIR.glob(f"{file_id}.*"))
            if not found_files: raise HTTPException(status_code=404, detail="Original file not found.")
            orig_path = found_files[0]; ext = orig_path.suffix.lower()

        # 2. Extract the specific image bytes & Context
        image_bytes = None
        target_shape = None
        doc = None
        pres = None
        ctx = {"surrounding_text": ""}
        
        if ext == ".pptx":
            pres = Presentation(orig_path)
            images = []
            valid_types = [
                MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE,
                MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.IGX_GRAPHIC, MSO_SHAPE_TYPE.GROUP
            ]
            for s_num, slide in enumerate(pres.slides, 1):
                for shape in slide.shapes:
                    if getattr(shape, "shape_type", None) in valid_types:
                        images.append((s_num, shape))
                        
            if 1 <= req.image_idx <= len(images):
                target_shape = images[req.image_idx - 1][1]
                # Native extraction for basic pictures
                if getattr(target_shape, "shape_type", None) in [MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE]:
                    try: 
                        if hasattr(target_shape, "image"): image_bytes = target_shape.image.blob
                    except AttributeError: pass
                # Note: Charts/Groups will have image_bytes=None here unless screenshot fallback is active
                
                from src.pipelines.agent_pipeline import get_context_for_image_pptx
                slide_target = pres.slides[req.slide_num - 1] if req.slide_num else pres.slides[images[req.image_idx-1][0]-1]
                ctx = get_context_for_image_pptx(slide_target, target_shape)
            else:
                raise HTTPException(status_code=404, detail="Image index out of bounds in PPTX.")
                
        elif ext == ".docx":
            doc = Document(orig_path)
            drawing_elements = doc._element.xpath('.//w:drawing')
            if 1 <= req.image_idx <= len(drawing_elements):
                target_shape = drawing_elements[req.image_idx - 1]
                rId = None
                blips = target_shape.xpath('.//a:blip/@r:embed')
                if blips: rId = blips[0]
                else:
                    charts = target_shape.xpath('.//c:chart/@r:id')
                    if charts: rId = charts[0]

                if rId:
                    rel = doc.part.rels[rId]
                    if "xml" not in rel.target_ref.lower():
                        image_bytes = rel.target_part.blob
                
                # Context extraction directly from XML
                parent_p = target_shape.xpath('./ancestor::w:p')
                if parent_p: ctx["surrounding_text"] = "".join(parent_p[0].itertext())[:300]
            else:
                raise HTTPException(status_code=404, detail="Image index out of bounds in DOCX.")

        # --- ADD THIS BLOCK: Recover missing image bytes from the JSON cache ---
        if not image_bytes:
            import base64
            json_path = PROCESSED_DIR / f"{file_id}.json"
            if json_path.exists():
                with open(json_path, "r", encoding="utf-8") as f:
                    cached_results = json.load(f)
                for res in cached_results:
                    if res.get("image_idx") == req.image_idx:
                        # Find the base64 string (often saved as 'base64_image' or 'image')
                        b64_str = res.get("base64_image") or res.get("image")
                        if b64_str:
                            if "," in b64_str: 
                                b64_str = b64_str.split(",")[1] # Strip data URI prefix if present
                            image_bytes = base64.b64decode(b64_str)
                            print(f"--> SUCCESS: Recovered missing image bytes from JSON cache!")
                            break
        # -----------------------------------------------------------------------

        # 3. Route to the correct pipeline
        from src.pipelines.agent_pipeline import classify_and_generate_alt_text, get_primary_model

        # 3. Route to the correct pipeline
        from src.pipelines.agent_pipeline import classify_and_generate_alt_text, get_primary_model
        primary_model = get_primary_model(provider=req.provider)
            
        _, gen_alt, _ = classify_and_generate_alt_text(
            image_bytes=image_bytes, 
            structured_context=ctx, 
            primary_model_system=primary_model, 
            ext=ext, 
            existing_alt="", 
            slide_num=req.slide_num,   
            task_idx=req.image_idx-1,
            provider=req.provider,     
            api_key=req.api_key,       
            forced_pipeline=req.forced_pipeline 
        )
        new_alt_text = gen_alt
        doc_modified = False

        # Apply the new alt text to the original image XML
        if ext == ".docx":
            docPr_elements = target_shape.xpath('.//wp:docPr')
            if docPr_elements:
                docPr_elements[0].set('descr', new_alt_text)
                doc_modified = True
        elif ext == ".pptx":
            cNvPr_elements = target_shape._element.xpath('.//*[local-name()="cNvPr"]')
            if cNvPr_elements:
                cNvPr_elements[0].set('descr', new_alt_text)
                doc_modified = True

        if req.forced_pipeline == "Equation" and image_bytes:
            from src.pipelines.equation_pipeline import extract_equations_from_image, insert_equation_into_docx
            equations = extract_equations_from_image(image_bytes, provider=req.provider, api_key=req.api_key)
            if equations:
                if ext == ".docx":
                    # Fallback insertion (since shape is raw XML now)
                    pass 
            else:
                # --- NEW: Bubble up Equation failure to the frontend ---
                raise HTTPException(
                    status_code=422, 
                    detail="Equation extraction failed: Could not identify mathematical content in the image."
                )
                
        elif req.forced_pipeline == "Table" and image_bytes:
            from src.pipelines.table_pipeline import extract_table_from_image, insert_table_into_docx, insert_table_into_pptx
            
            logging.info("--> ENTERED TABLE PIPELINE ROUTING")
            table_data = extract_table_from_image(image_bytes, provider=req.provider, api_key=req.api_key)
            logging.info(f"--> EXTRACTED TABLE DATA: {table_data}") 
            
            if table_data and any(table_data):
                if ext == ".docx":
                    insert_table_into_docx(doc, target_shape, table_data, alt_text=new_alt_text)
                    doc_modified = True
                    logging.info("--> DOCX TABLE INSERTED")
                elif ext == ".pptx":
                    insert_table_into_pptx(slide_target, target_shape, table_data, alt_text=new_alt_text)
                    doc_modified = True
                    logging.info("--> PPTX TABLE INSERTED")
            else:
                # --- NEW: Bubble up Table failure to the frontend ---
                logging.warning("--> TABLE DATA WAS EMPTY, ABORTING AND NOTIFYING FRONTEND")
                raise HTTPException(
                    status_code=422, 
                    detail="Table extraction failed: Could not detect a valid table structure in the image."
                )
        if doc_modified:
            if ext == ".docx" and doc:
                doc.save(orig_path)
                logging.info(f"--> Saved modified DOCX back to {orig_path}")
            elif ext == ".pptx" and pres:
                pres.save(orig_path)
                logging.info(f"--> Saved modified PPTX back to {orig_path}")
        # 5. Update the JSON processing cache
        json_path = PROCESSED_DIR / f"{file_id}.json"
        if json_path.exists():
            with open(json_path, "r", encoding="utf-8") as f:
                results = json.load(f)
            for res in results:
                if res.get("image_idx") == req.image_idx:
                    if req.forced_pipeline in res["classification"]:
                        res["classification"].remove(req.forced_pipeline)
                    res["classification"].insert(0, req.forced_pipeline)
                    res["generated_alt_text"] = new_alt_text
                    res["final_alt_text"] = new_alt_text
                    break
            with open(json_path, "w", encoding="utf-8") as f:
                json.dump(results, f)

        return {"status": "success", "new_alt_text": new_alt_text, "pipeline_used": req.forced_pipeline}

    except Exception as e:
        logging.error(f"Failed to regenerate image {req.image_idx} for {file_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))

@app.get("/download/{file_id}")
def download_file(file_id: str):
    """Applies final alt text and returns the remediated file."""
    ext = None; orig_path = None; orig_filename = "unknown_file"
    try:
        if file_id in processing_status and "filename" in processing_status[file_id]:
            orig_filename = processing_status[file_id]["filename"]
            potential_ext = os.path.splitext(orig_filename)[1].lower()
            if potential_ext in [".docx", ".pptx"]:
                ext = potential_ext
                orig_path_check = UPLOAD_DIR / f"{file_id}{ext}"
                if orig_path_check.exists(): orig_path = orig_path_check
                else: 
                     found_files = list(UPLOAD_DIR.glob(f"{file_id}.*"))
                     if found_files: orig_path = found_files[0]; ext = orig_path.suffix.lower()
                     else: raise HTTPException(status_code=404, detail="Original uploaded file not found.")
            else: raise HTTPException(status_code=400, detail="Invalid extension.")
        else:
            found_files = list(UPLOAD_DIR.glob(f"{file_id}.*"))
            if found_files:
                orig_path = found_files[0]; ext = orig_path.suffix.lower(); orig_filename = orig_path.name 
            else: raise HTTPException(status_code=404, detail="Original uploaded file not found.")

        if ext not in [".docx", ".pptx"]: raise HTTPException(status_code=400, detail=f"Unsupported file type: {ext}")
    except Exception as e: raise HTTPException(status_code=500, detail=f"Error finding original file: {e}")

    json_path = PROCESSED_DIR / f"{file_id}.json"
    if not json_path.exists(): raise HTTPException(status_code=404, detail="Alt text processing data not found.")
    try:
        with open(json_path, "r", encoding="utf-8") as f: results = json.load(f)
    except Exception as e: raise HTTPException(status_code=500, detail="Error reading alt text data.")

    out_path = REMEDIATED_DIR / f"remediated_{file_id}{ext}"
    safe_orig_filename = "".join(c for c in orig_filename if c.isalnum() or c in (' ', '.', '_')).rstrip()
    download_filename = f"remediated_{safe_orig_filename}"

    try:
        if ext == ".docx":
            import zipfile
            from lxml import etree
            import tempfile
            
            temp_dir = tempfile.mkdtemp()
            try:
                # 1. Unzip the docx exactly like the notebook
                with zipfile.ZipFile(orig_path, 'r') as docx_zip:
                    docx_zip.extractall(temp_dir)

                document_xml_path = os.path.join(temp_dir, 'word', 'document.xml')
                tree = etree.parse(document_xml_path)
                root = tree.getroot()

                
                namespaces = {
                    'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                    'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing'
                }
                
                # 2. Iterate through drawings and apply alt text
                drawings = root.xpath('.//w:drawing', namespaces=namespaces) 
                
                for i, drawing in enumerate(drawings):
                    if i < len(results):
                        final_alt_text = results[i].get("final_alt_text", results[i].get("generated_alt_text", ""))
                        
                        # MUST include namespaces=namespaces here for lxml
                        docPr_elements = drawing.xpath('.//wp:docPr', namespaces=namespaces) 
                        
                        if docPr_elements:
                            docPr_elements[0].set('descr', final_alt_text)

                # 3. Save modified XML
                tree.write(document_xml_path, encoding='UTF-8', xml_declaration=True)

                # 4. Re-zip into the remediated directory
                with zipfile.ZipFile(out_path, 'w') as outzip:
                    for root_dir, _, files in os.walk(temp_dir):
                        for file in files:
                            file_path = os.path.join(root_dir, file)
                            arcname = os.path.relpath(file_path, temp_dir)
                            outzip.write(file_path, arcname)
            finally:
                import shutil
                shutil.rmtree(temp_dir)

        elif ext == ".pptx":
            pres = Presentation(orig_path)
            valid_types = [
                MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE,
                MSO_SHAPE_TYPE.CHART, MSO_SHAPE_TYPE.IGX_GRAPHIC, MSO_SHAPE_TYPE.GROUP
            ]
            shapes_processed = 0
            for slide in pres.slides:
                for shape in slide.shapes:
                    if getattr(shape, "shape_type", None) in valid_types:
                         if shapes_processed < len(results):
                            final_alt_text = results[shapes_processed].get("final_alt_text", results[shapes_processed].get("generated_alt_text", ""))
                            try:
                                cNvPr_elements = shape._element.xpath('.//*[local-name()="cNvPr"]')
                                if cNvPr_elements: cNvPr_elements[0].set('descr', final_alt_text)
                            except Exception: pass
                            shapes_processed += 1
            pres.save(out_path)

    except Exception as e:
         logging.exception(f"Failed to remediate file {file_id}: {e}") 
         raise HTTPException(status_code=500, detail=f"Failed to write remediated file: {str(e)}")

    return FileResponse(
        out_path, filename=download_filename,
        headers={"Content-Disposition": f"attachment; filename*=UTF-8''{download_filename}"}
    )

if __name__ == "__main__":
    import uvicorn
    # --- START MODIFICATION: Ensure necessary directories exist before starting ---
    UPLOAD_DIR.mkdir(parents=True, exist_ok=True)
    PROCESSED_DIR.mkdir(parents=True, exist_ok=True)
    REMEDIATED_DIR.mkdir(parents=True, exist_ok=True)
    CUSTOM_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    # --- END MODIFICATION ---

    logging.basicConfig(level=logging.INFO)
    logging.info(f"Starting server... Uploads: {UPLOAD_DIR}, Processed: {PROCESSED_DIR}, Remediated: {REMEDIATED_DIR}, Cache: {CUSTOM_CACHE_DIR}")
    uvicorn.run("main:app", host="0.0.0.0", port=8000, reload=True)
